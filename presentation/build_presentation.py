from __future__ import annotations

from io import BytesIO
from pathlib import Path
import shutil
import subprocess
import zipfile

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "presentation" / "assets"
DESTINATION = ROOT / "presentation" / "Event-TimeRAF_Verified_Presentation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

INK = RGBColor(29, 38, 48)
MUTED = RGBColor(86, 99, 112)
TEAL = RGBColor(31, 123, 117)
TEAL_LIGHT = RGBColor(225, 242, 239)
GOLD = RGBColor(192, 133, 35)
GOLD_LIGHT = RGBColor(250, 240, 213)
ROSE = RGBColor(174, 74, 91)
ROSE_LIGHT = RGBColor(248, 228, 231)
BLUE_LIGHT = RGBColor(227, 238, 248)
GREEN_LIGHT = RGBColor(229, 243, 232)
PAPER = RGBColor(250, 251, 252)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(205, 213, 220)


def add_text(slide, text, x, y, w, h, size=20, bold=False, color=INK,
             align=PP_ALIGN.LEFT, font="Aptos", margin=0.04,
             valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rich_lines(slide, lines, x, y, w, h, size=18, gap=7):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    for index, (lead, body, color) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(gap)
        run = paragraph.add_run()
        run.text = lead
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = color
        run = paragraph.add_run()
        run.text = body
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.color.rgb = INK
    return box


def add_panel(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def add_title(slide, title, subtitle=None, number=None):
    title_size = 22 if len(title) > 42 else 25
    add_text(slide, title, 0.55, 0.30, 11.8, 0.58, title_size, True, INK)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.83, 11.7, 0.32, 11, False, MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.14), Inches(12.15), Inches(0.025))
    line.fill.solid(); line.fill.fore_color.rgb = LINE; line.line.fill.background()
    if number is not None:
        add_text(slide, f"{number:02d}", 12.45, 0.35, 0.4, 0.32, 10, True, TEAL, PP_ALIGN.RIGHT)


def add_footer(slide, text="Verified against immutable run 20260723T112033170131Z"):
    add_text(slide, text, 0.62, 7.15, 11.8, 0.20, 8, False, MUTED)


def add_picture_contain(slide, path: Path, x, y, w, h):
    from PIL import Image
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    return slide.shapes.add_picture(
        str(path), Inches(x + (w - dw) / 2), Inches(y + (h - dh) / 2),
        width=Inches(dw), height=Inches(dh),
    )


def add_table(slide, data, x, y, w, h, widths=None, font_size=11,
              header_fill=INK, highlight_rows=None):
    rows = len(data)
    cols = len(data[0])
    shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    if widths:
        for index, width in enumerate(widths):
            table.columns[index].width = Inches(width)
    for row_index, row in enumerate(data):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill = header_fill if row_index == 0 else WHITE
            if highlight_rows and row_index in highlight_rows:
                fill = highlight_rows[row_index]
            cell.fill.solid(); cell.fill.fore_color.rgb = fill
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT if col_index == 0 else PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(font_size)
                    run.font.bold = row_index == 0 or (highlight_rows and row_index in highlight_rows)
                    run.font.color.rgb = WHITE if row_index == 0 else INK
    return table


def notes(slide, text):
    frame = slide.notes_slide.notes_text_frame
    frame.text = text


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid(); background.fore_color.rgb = PAPER
    return slide


def prepare_assets():
    ASSETS.mkdir(parents=True, exist_ok=True)
    methodology = [
        "event_timeraf_pipeline_overview",
        "source_audited_context_builder",
        "leakage_safe_event_context_retriever",
        "drift_aware_forecast_evidence_head",
    ]
    for stem in methodology:
        subprocess.run([
            "pdftoppm", "-png", "-singlefile", "-r", "160",
            str(ROOT / "paper" / "figures" / f"{stem}.pdf"),
            str(ASSETS / stem),
        ], check=True)
    shutil.copy2(ROOT / "paper" / "figures" / "forecast_case.png", ASSETS / "forecast_case.png")

    archive = ROOT / "event_timeraf_final_run.zip"
    with zipfile.ZipFile(archive) as bundle:
        results = pd.read_csv(BytesIO(bundle.read("outputs/tables/main_results.csv")))
    results = results.loc[results["subset"].eq("all")].copy()
    order = [
        "M00_persistence", "M01_daily_seasonal", "M02_weekly_seasonal",
        "M03_xgb_pm25", "M04_xgb_context", "M05_random_retrieval",
        "M06_cosine_retrieval", "M07_xgb_cosine",
        "M08_event_timeraf_no_drift", "M09_event_timeraf_full",
        "M10_frozen_chronos", "M11_chronos_hybrid_retrieval",
    ]
    labels = [item.split("_", 1)[0] for item in order]
    values = results.set_index("model").loc[order, "mse"].to_numpy()
    colors = ["#98A6B3"] * len(order)
    colors[4] = "#1F7B75"
    colors[9] = "#C08523"
    colors[10] = colors[11] = "#7099B8"
    figure, axis = plt.subplots(figsize=(9.6, 5.2))
    axis.barh(labels[::-1], values[::-1], color=colors[::-1], height=0.68)
    axis.set_xlabel("Test MSE (lower is better)")
    axis.set_xlim(0, 70)
    axis.grid(axis="x", alpha=0.20)
    axis.spines[["top", "right", "left"]].set_visible(False)
    for index, value in enumerate(values[::-1]):
        axis.text(value + 0.7, index, f"{value:.3f}", va="center", fontsize=9)
    figure.tight_layout()
    figure.savefig(ASSETS / "main_mse_bars.png", dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(figure)


def build():
    prepare_assets()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "Event-TimeRAF Verified Presentation"
    prs.core_properties.subject = "Manifest-backed PM2.5 forecasting study"

    # 1. Title
    slide = blank_slide(prs)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.20), SLIDE_H).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = TEAL; slide.shapes[-1].line.fill.background()
    add_text(slide, "EVENT-TIMERAF", 0.75, 1.25, 5.8, 0.45, 15, True, TEAL)
    add_text(slide, "Verified event-aware retrieval for\n24-hour PM2.5 forecasting", 0.72, 1.75, 7.0, 1.42, 34, True, INK)
    add_text(slide, "Los Angeles County | 2019-2024 | Faculty research update", 0.76, 3.45, 6.8, 0.38, 15, False, MUTED)
    add_panel(slide, 8.15, 1.20, 4.15, 4.50, TEAL_LIGHT, TEAL)
    add_text(slide, "VERIFIED SCOPE", 8.55, 1.62, 3.3, 0.34, 13, True, TEAL)
    add_rich_lines(slide, [
        ("Input  ", "168 hours", TEAL),
        ("Output  ", "24 hourly PM2.5 values", TEAL),
        ("Test  ", "7,199 forecast origins", TEAL),
        ("Best  ", "M04 context XGBoost", GOLD),
        ("Status  ", "Pass with limitations", ROSE),
    ], 8.48, 2.15, 3.45, 2.65, 17, 8)
    add_text(slide, "github.com/shasabbir/event-time-raf", 0.76, 6.78, 6.5, 0.25, 10, False, MUTED)
    notes(slide, "Open with the problem and verified scope. State immediately that M04, not the full M09 model, is the strongest manifest-backed result.")

    # 2. Question and evidence contract
    slide = blank_slide(prs); add_title(slide, "Research question and evidence contract", "What the experiment tests, and what it does not claim", 2)
    add_panel(slide, 0.62, 1.45, 5.85, 5.35, WHITE, LINE)
    add_text(slide, "Research question", 0.94, 1.76, 4.9, 0.35, 17, True, TEAL)
    add_text(slide, "Can source-audited weather, calendar, and storm-event context improve leakage-safe retrieval and forecasting under distribution shift?", 0.94, 2.20, 5.15, 1.55, 21, True, INK)
    add_rich_lines(slide, [
        ("Input example:  ", "7 days of hourly PM2.5 + aligned context", TEAL),
        ("Target:  ", "the next 24 hourly PM2.5 values", TEAL),
        ("Evaluation:  ", "MSE, MAE, RMSE, R² + paired bootstrap", TEAL),
    ], 0.92, 3.66, 5.12, 1.65, 15, 6)
    add_panel(slide, 6.78, 1.45, 5.93, 5.35, BLUE_LIGHT, TEAL)
    add_text(slide, "Evidence contract", 7.10, 1.76, 4.9, 0.35, 17, True, TEAL)
    add_rich_lines(slide, [
        ("69/69  ", "manifest SHA-256 checks pass", TEAL),
        ("M00-M11  ", "saved predictions reproduced", TEAL),
        ("M12  ", "excluded; created after the manifest", ROSE),
        ("Events  ", "retrospective, not strict real time", ROSE),
        ("Geography  ", "one county; no external validation", ROSE),
    ], 7.05, 2.25, 5.20, 2.85, 18, 9)
    add_text(slide, "No unsupported result remains in the corrected paper.", 7.10, 5.70, 4.9, 0.55, 18, True, INK)
    add_footer(slide)
    notes(slide, "Explain the immutable ZIP as the source of truth. The notebook does not silently choose between conflicting edited artifacts.")

    # 3. Pipeline
    slide = blank_slide(prs); add_title(slide, "End-to-end framework", "SACB -> LSER -> DFEH, with an auditable evidence record", 3)
    add_picture_contain(slide, ASSETS / "event_timeraf_pipeline_overview.png", 0.55, 1.25, 12.25, 5.65)
    add_footer(slide)
    notes(slide, "Walk left to right. Official sources enter SACB, only train-history candidates enter LSER, and DFEH evaluates direct XGBoost and frozen Chronos paths. Point out that the drift gate is marked unreported.")

    # 4. SACB
    slide = blank_slide(prs); add_title(slide, "Module 1: Source-Audited Context Builder", "Source checks, causal alignment, and chronological windows", 4)
    add_panel(slide, 0.62, 1.42, 4.00, 5.45, TEAL_LIGHT, TEAL)
    add_text(slide, "Origin context", 0.95, 1.78, 3.1, 0.35, 17, True, TEAL)
    add_text(slide, "qₜ ∈ ℝ⁸⁵", 0.94, 2.17, 3.1, 0.48, 27, True, INK)
    add_table(slide, [
        ["Feature group", "Count"], ["PM2.5 history", "23"], ["Weather", "14"],
        ["Calendar", "9"], ["Storm event", "39"], ["Total", "85"],
    ], 0.93, 2.85, 3.37, 2.35, [2.45, 0.92], 12, highlight_rows={5: GOLD_LIGHT})
    add_text(slide, "Target values are never interpolated. Filling uses past observations only.", 0.94, 5.48, 3.2, 0.90, 13, True, ROSE)
    add_picture_contain(slide, ASSETS / "source_audited_context_builder.png", 4.85, 1.36, 7.85, 5.55)
    add_footer(slide)
    notes(slide, "Derive 85 from memory: 23 PM2.5, 14 weather, 9 calendar, 39 event. Explain that 86/40 was an AI draft counting error corrected from actual feature names.")

    # 5. LSER
    slide = blank_slide(prs); add_title(slide, "Module 2: Leakage-Safe Event-Context Retriever", "Temporal eligibility before similarity ranking", 5)
    add_picture_contain(slide, ASSETS / "leakage_safe_event_context_retriever.png", 0.52, 1.30, 8.15, 5.75)
    add_panel(slide, 8.82, 1.46, 3.88, 5.25, GOLD_LIGHT, GOLD)
    add_text(slide, "Embargo", 9.15, 1.78, 3.1, 0.35, 17, True, GOLD)
    add_text(slide, "candidate target end\n< query input start", 9.14, 2.20, 3.1, 0.85, 20, True, INK)
    add_text(slide, "Hybrid score", 9.15, 3.42, 3.1, 0.35, 17, True, GOLD)
    add_text(slide, "0.5 shape + 0.2 weather\n+ 0.1 calendar + 0.2 event", 9.14, 3.83, 3.15, 0.88, 18, False, INK)
    add_rich_lines(slide, [
        ("KB  ", "191 non-overlapping train windows", GOLD),
        ("Top-k  ", "k = 8; sensitivity {1,4,8,16}", GOLD),
    ], 9.08, 5.15, 3.2, 1.0, 14, 5)
    add_footer(slide)
    notes(slide, "The leakage rule is stricter than candidate origin before query origin. The entire candidate future must end before the query lookback starts.")

    # 6. DFEH
    slide = blank_slide(prs); add_title(slide, "Module 3: Drift-Aware Forecast and Evidence Head", "Two verified forecast paths; one excluded selector", 6)
    add_picture_contain(slide, ASSETS / "drift_aware_forecast_evidence_head.png", 0.48, 1.28, 8.40, 5.80)
    add_panel(slide, 9.00, 1.44, 3.70, 5.28, ROSE_LIGHT, ROSE)
    add_text(slide, "Direct path", 9.32, 1.78, 2.9, 0.32, 16, True, ROSE)
    add_text(slide, "24 XGBoost regressors\n151 inputs per horizon", 9.31, 2.14, 3.0, 0.72, 18, True, INK)
    add_text(slide, "Frozen path", 9.32, 3.18, 2.9, 0.32, 16, True, ROSE)
    add_text(slide, "Chronos-Bolt + retrieval\nω selected on validation", 9.31, 3.55, 3.0, 0.72, 17, False, INK)
    add_text(slide, "Evidence output", 9.32, 4.62, 2.9, 0.32, 16, True, ROSE)
    add_text(slide, "feature effects, retrieved IDs, event IDs, drift, uncertainty", 9.31, 4.99, 3.0, 0.90, 13, False, INK)
    add_text(slide, "Drift gate: excluded until a complete manifest-backed rerun.", 9.30, 5.96, 3.0, 0.58, 10, True, ROSE)
    add_footer(slide)
    notes(slide, "Derive 151: 85 context + 51 retrieval summary + 6 drift fields + 9 future-calendar fields. The five drift components are robust-scaled on training data; the threshold is the validation 0.90 quantile.")

    # 7. Protocol
    slide = blank_slide(prs); add_title(slide, "Experimental protocol", "Chronological evaluation and pre-specified model ladder", 7)
    add_panel(slide, 0.62, 1.43, 4.02, 5.28, WHITE, LINE)
    add_text(slide, "Data split", 0.95, 1.77, 3.2, 0.35, 17, True, TEAL)
    add_text(slide, "34,020", 0.95, 2.25, 1.35, 0.45, 24, True, INK)
    add_text(slide, "TRAIN", 2.32, 2.34, 1.3, 0.25, 11, True, MUTED)
    add_text(slide, "7,113", 0.95, 2.95, 1.35, 0.45, 24, True, INK)
    add_text(slide, "VALIDATION", 2.32, 3.04, 1.5, 0.25, 11, True, MUTED)
    add_text(slide, "7,199", 0.95, 3.65, 1.35, 0.45, 24, True, INK)
    add_text(slide, "TEST", 2.32, 3.74, 1.3, 0.25, 11, True, MUTED)
    add_text(slide, "172,776 test points", 0.95, 4.48, 2.8, 0.35, 17, True, GOLD)
    add_text(slide, "500 paired block-bootstrap resamples\n24-hour blocks", 0.95, 5.15, 3.1, 0.78, 15, False, INK)
    add_table(slide, [
        ["IDs", "Purpose"],
        ["M00-M02", "Persistence + seasonal"],
        ["M03-M04", "XGBoost PM2.5/context"],
        ["M05-M06", "Random/cosine retrieval"],
        ["M07", "XGBoost + cosine"],
        ["M08-M09", "Event-TimeRAF no/full drift"],
        ["M10-M11", "Frozen Chronos/fusion"],
    ], 4.92, 1.48, 7.78, 4.02, [1.58, 6.20], 14, highlight_rows={3: GOLD_LIGHT, 5: TEAL_LIGHT, 6: BLUE_LIGHT})
    add_text(slide, "Metrics", 4.98, 5.78, 1.2, 0.30, 15, True, TEAL)
    add_text(slide, "MSE: large-error penalty | MAE: absolute error | RMSE: PM2.5 units | R²: explained variation", 6.00, 5.72, 6.25, 0.75, 13, False, INK)
    add_footer(slide)
    notes(slide, "Explain that model choices and validation decisions precede test examination. Difference A minus B below zero favors A.")

    # 8. Main results
    slide = blank_slide(prs); add_title(slide, "Verified test results", "The strongest verified model is the context baseline M04", 8)
    add_picture_contain(slide, ASSETS / "main_mse_bars.png", 0.50, 1.33, 7.65, 5.65)
    add_panel(slide, 8.42, 1.52, 4.28, 5.10, TEAL_LIGHT, TEAL)
    add_text(slide, "Key results", 8.78, 1.85, 3.35, 0.34, 17, True, TEAL)
    add_table(slide, [
        ["Model", "MSE", "MAE"],
        ["M04 context", "26.185", "3.125"],
        ["M09 full", "26.712", "3.149"],
        ["M10 Chronos", "28.941", "3.205"],
        ["M11 fusion", "28.709", "3.209"],
    ], 8.77, 2.36, 3.55, 2.20, [1.75, 0.90, 0.90], 12, highlight_rows={1: GREEN_LIGHT})
    add_text(slide, "Interpretation", 8.78, 4.92, 3.3, 0.30, 16, True, GOLD)
    add_text(slide, "M09 does not beat M04 overall. M11 improves M10 MSE slightly, but not conclusively.", 8.77, 5.30, 3.45, 1.05, 15, True, INK)
    add_footer(slide)
    notes(slide, "Do not say the proposed full model wins. Say M04 is strongest, Event-TimeRAF is an audited framework, and retrieval gains are selective.")

    # 9. Ablation
    slide = blank_slide(prs); add_title(slide, "Ablation study, run for real", "Paired MSE differences; negative values favor the first model", 9)
    add_table(slide, [
        ["Comparison", "ΔMSE", "95% interval", "Conclusion"],
        ["M04 - M03: weather/calendar", "-3.764", "[-5.856, -2.071]", "clear gain"],
        ["M06 - M05: cosine/random", "-4.284", "[-5.862, -2.759]", "clear gain"],
        ["M07 - M04: add retrieval", "+0.487", "[0.015, 1.028]", "worse"],
        ["M08 - M07: event ranking", "+0.005", "[-0.292, 0.275]", "no clear gain"],
        ["M09 - M08: drift fields", "+0.035", "[-0.101, 0.161]", "no clear gain"],
        ["M09 - A00: event context", "+0.083", "[-0.225, 0.354]", "no clear gain"],
        ["M11 - M10: Chronos fusion", "-0.233", "[-0.665, 0.152]", "favorable, uncertain"],
    ], 0.64, 1.42, 12.05, 4.65, [4.40, 1.15, 2.55, 3.95], 13,
        highlight_rows={1: GREEN_LIGHT, 2: GREEN_LIGHT, 3: ROSE_LIGHT})
    add_panel(slide, 0.68, 6.25, 12.00, 0.58, GOLD_LIGHT, GOLD)
    add_text(slide, "Justified: audited context and cosine ranking. Not justified: a uniform gain from event, drift, or retrieval features on top of M04.", 0.95, 6.33, 11.45, 0.40, 13, True, INK)
    add_footer(slide)
    notes(slide, "Every row is calculated from actual saved predictions. Explain confidence intervals: a crossing of zero means the direction is not conclusive under this bootstrap.")

    # 10. Evidence and limitations
    slide = blank_slide(prs); add_title(slide, "Explainability and scope limits", "Traceable records without overclaiming operational readiness", 10)
    add_picture_contain(slide, ASSETS / "forecast_case.png", 0.55, 1.38, 6.30, 3.55)
    add_panel(slide, 0.70, 5.10, 6.00, 1.45, BLUE_LIGHT, TEAL)
    add_text(slide, "7,199 evidence records", 1.00, 5.36, 4.85, 0.38, 17, True, TEAL)
    add_text(slide, "Top XGBoost effects | retrieved IDs | event IDs | drift score | uncertainty proxy", 1.00, 5.84, 5.25, 0.48, 11, False, INK)
    add_panel(slide, 7.10, 1.38, 5.60, 5.18, ROSE_LIGHT, ROSE)
    add_text(slide, "Claims that remain limited", 7.46, 1.75, 4.55, 0.36, 18, True, ROSE)
    add_rich_lines(slide, [
        ("Event timing  ", "no publication timestamps", ROSE),
        ("External validity  ", "one county only", ROSE),
        ("Random replay  ", "requires NumPy 2.0.2 or saved evidence", ROSE),
        ("Efficiency  ", "hardware identifiers were not recorded", ROSE),
        ("TimeRAF scope  ", "no learned dual encoder or Channel Prompting", ROSE),
    ], 7.40, 2.42, 4.80, 2.95, 16, 9)
    add_text(slide, "These are reported limitations, not hidden failures.", 7.45, 5.74, 4.55, 0.50, 17, True, INK)
    add_footer(slide)
    notes(slide, "Explainability is a structured evidence record, not a causal explanation. Events are official but retrospective because issue times are unavailable.")

    # 11. Verification and live run
    slide = blank_slide(prs); add_title(slide, "What changed after verification", "Live notebook: notebooks/03_paper_claim_verification.ipynb", 11)
    add_table(slide, [
        ["AI draft", "Verified correction"],
        ["86 context / 40 event / 152 input", "85 context / 39 event / 151 input"],
        ["M12 as the best final result", "M12 removed: not manifest-backed"],
        ["Generic drift equation", "Five statistics + median/MAD + validation threshold"],
        ["Operational event implication", "Retrospective event sensitivity only"],
    ], 0.65, 1.42, 7.15, 3.45, [3.32, 3.83], 14,
        highlight_rows={1: GOLD_LIGHT, 2: ROSE_LIGHT})
    add_panel(slide, 8.10, 1.42, 4.58, 4.95, TEAL_LIGHT, TEAL)
    add_text(slide, "Live run order", 8.46, 1.78, 3.7, 0.34, 18, True, TEAL)
    add_rich_lines(slide, [
        ("1  ", "Verify 69 hashes", TEAL),
        ("2  ", "Rebuild data and retrieval", TEAL),
        ("3  ", "Rerun M00-M11", TEAL),
        ("4  ", "Recompute metrics + ablations", TEAL),
        ("5  ", "Regenerate nine figures", TEAL),
        ("6  ", "Pass final claim gate", TEAL),
    ], 8.39, 2.38, 3.75, 2.90, 17, 6)
    add_text(slide, "Expected runtime: about 2-3 minutes with the archived ZIP attached.", 8.45, 5.32, 3.65, 0.78, 12, True, INK)
    add_text(slide, "Bottom line", 0.78, 5.28, 2.0, 0.30, 16, True, GOLD)
    add_text(slide, "The study is successful as a reproducible, leakage-audited framework with selective retrieval gains, not as a claim that Event-TimeRAF beats every baseline.", 0.76, 5.68, 6.60, 1.00, 16, True, INK)
    add_footer(slide, "Notebook, paper, verification log, slides, and source ZIP use the same archived run")
    notes(slide, "Finish with the qualified conclusion. Then open the notebook and run all cells. Keep the immutable ZIP in the project root or set FINAL_RUN_ZIP_OVERRIDE.")

    DESTINATION.parent.mkdir(parents=True, exist_ok=True)
    prs.save(DESTINATION)
    print(DESTINATION)


if __name__ == "__main__":
    build()
